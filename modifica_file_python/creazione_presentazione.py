from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Schema Riassuntivo – Pipeline Terraform per Uffici"
subtitle.text = "Sintesi della nuova architettura e dei processi operativi"

content = [
    ("Obiettivo", [
        "Pipeline per ufficio, non per risorsa",
        "Ogni ufficio gestisce moduli Terraform con statefile dedicato",
        "Nuovi moduli → si assegnano all’ufficio, nessuna nuova pipeline"
    ]),
    ("Situazione Attuale", [
        "Pipeline legate a progetto/ambiente",
        "Unico statefile per progetto",
        "Codice condiviso tra uffici"
    ]),
    ("Nuovo Modello Operativo", [
        "1 pipeline per ufficio",
        "Ogni pipeline esegue solo i moduli di competenza",
        "Esempi: CTDUMA-SYSOPS, CTDUMA-ASA, CTDUMA-CLOUD, COMMON"
    ]),
    ("Statefile", [
        "Separati per ufficio/ambiente/modulo",
        "Esempio: sysops/dev/infrastructure.tfstate"
    ]),
    ("Permessi ADO", [
        "ASA non può lanciare CTDUMA-SYSOPS",
        "SYSOPS non può lanciare CTDUMA-ASA",
        "Codice modificabile da tutti"
    ]),
    ("Moduli Common", [
        "Gestiti in pipeline COMMON",
        "Per risorse non associate a un ufficio specifico"
    ]),
    ("Passi Operativi", [
        "Mappatura moduli → assegnazione uffici",
        "Creazione pipeline ufficio",
        "Config statefile dedicati",
        "Aggiornamento permessi ADO",
        "Validazione con terraform plan"
    ]),
    ("Benefici", [
        "Governance chiara",
        "Scalabilità",
        "Sicurezza tramite permessi ADO",
        "Indipendenza tramite statefile separati"
    ]),
]

# Generate slides
for title_text, bullets in content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = title_text
    tf = body.text_frame
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 1

file_path = "/mnt/data/Schema_Pipeline_Terraform_Uffici.pptx"
prs.save(file_path)

file_path
